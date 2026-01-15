# file_utils.py

import os
import io
import re
import tempfile
import fitz  # PyMuPDF
import pandas as pd
import win32com.client  # Pour les .doc (Windows seulement)

from typing import List, Tuple
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

# Gestion des variables d'environnement
from dotenv import load_dotenv
from azure.ai.vision.imageanalysis import ImageAnalysisClient
from azure.ai.vision.imageanalysis.models import VisualFeatures
from azure.core.credentials import AzureKeyCredential


# Charger .env
load_dotenv()
AZURE_VISION_ENDPOINT = os.getenv("AZURE_VISION_ENDPOINT")
AZURE_VISION_KEY = os.getenv("AZURE_VISION_KEY")


# ----------- PDF -----------
def extract_text_from_pdf(file_path: str) -> str:
    """Extrait le texte d'un fichier PDF"""
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()
    return text


# ----------- TXT / MD -----------
def extract_text_from_txt(file_path: str) -> str:
    """Extrait le texte d'un fichier TXT ou MD"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(file_path, "r", encoding="latin-1") as f:
            return f.read()


# ----------- DOC (Windows uniquement) -----------
def extract_text_from_doc_windows(doc_bytes: bytes) -> str:
    """Extrait le texte d'un fichier DOC (Windows uniquement)"""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as tmp:
        tmp.write(doc_bytes)
        tmp_path = tmp.name

    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(tmp_path)
        text = doc.Content.Text
        doc.Close()
        word.Quit()
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
    return text


# ----------- DOCX -----------
def extract_text_from_docx_bytes(doc_bytes: bytes) -> str:
    """Extrait le texte d'un fichier DOCX"""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
        tmp.write(doc_bytes)
        tmp_path = tmp.name
    
    try:
        doc = Document(tmp_path)
        text = "\n".join([p.text for p in doc.paragraphs])
        
        # Extraire aussi le texte des tableaux
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += "\n" + cell.text
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
    
    return text


# ----------- CSV -----------
def extract_text_from_csv_bytes(csv_bytes: bytes) -> str:
    """Extrait le texte d'un fichier CSV"""
    try:
        df = pd.read_csv(io.BytesIO(csv_bytes))
        return df.to_string(index=False)
    except Exception as e:
        # Fallback: lire comme texte brut
        try:
            return csv_bytes.decode('utf-8')
        except UnicodeDecodeError:
            return csv_bytes.decode('latin-1')


# ----------- XLSX -----------
def extract_text_from_xlsx_bytes(xlsx_bytes: bytes) -> str:
    """Extrait le texte d'un fichier XLSX"""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
        tmp.write(xlsx_bytes)
        tmp_path = tmp.name
    
    try:
        wb = load_workbook(tmp_path)
        text = ""
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text += f"\n\n--- Feuille: {sheet_name} ---\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    text += row_text + "\n"
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
    
    return text.strip()


# ----------- PPTX -----------
def extract_text_from_pptx_bytes(pptx_bytes: bytes) -> str:
    """Extrait le texte d'un fichier PPTX"""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(pptx_bytes)
        tmp_path = tmp.name
    
    try:
        prs = Presentation(tmp_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"\n\n--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
    
    return text.strip()


# ----------- IMAGES (Azure Vision OCR) -----------
def extract_text_from_image_bytes_azure(image_bytes: bytes) -> str:
    """Extrait le texte d'une image via Azure Vision OCR"""
    if not AZURE_VISION_ENDPOINT or not AZURE_VISION_KEY:
        raise RuntimeError("⚠️ Azure Vision endpoint/clé manquants dans .env")

    try:
        client = ImageAnalysisClient(
            endpoint=AZURE_VISION_ENDPOINT,
            credential=AzureKeyCredential(AZURE_VISION_KEY)
        )

        result = client.analyze(
            image_data=image_bytes,
            visual_features=[VisualFeatures.READ]
        )

        text = []
        if result.read is not None:
            for block in result.read.blocks:
                for line in block.lines:
                    text.append(line.text)

        return "\n".join(text).strip()
    
    except Exception as e:
        raise RuntimeError(f"Erreur Azure Vision OCR: {str(e)}")


# ----------- Fonction principale -----------
async def extract_text(file) -> str:
    """Fonction principale pour extraire le texte de différents types de fichiers"""
    if hasattr(file, "file") and hasattr(file, "filename"):
        # Cas: fichier uploadé (FastAPI UploadFile)
        contents = await file.read()
        ext = os.path.splitext(file.filename)[1].lower()

        if ext == ".pdf":
            with fitz.open(stream=io.BytesIO(contents), filetype="pdf") as doc:
                return "".join([page.get_text() for page in doc])

        elif ext in [".txt", ".md"]:
            try:
                return contents.decode("utf-8")
            except UnicodeDecodeError:
                return contents.decode("latin-1")

        elif ext == ".docx":
            return extract_text_from_docx_bytes(contents)

        elif ext == ".doc":
            return extract_text_from_doc_windows(contents)

        elif ext == ".csv":
            return extract_text_from_csv_bytes(contents)

        elif ext == ".xlsx":
            return extract_text_from_xlsx_bytes(contents)

        elif ext == ".pptx":
            return extract_text_from_pptx_bytes(contents)

        elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif"]:
            return extract_text_from_image_bytes_azure(contents)

        else:
            raise ValueError(f"Format de fichier non supporté : {ext}")

    else:
        # Cas: chemin de fichier local
        ext = os.path.splitext(file)[1].lower()
        if ext == ".pdf":
            return extract_text_from_pdf(file)
        elif ext in [".txt", ".md"]:
            return extract_text_from_txt(file)
        elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif"]:
            with open(file, "rb") as f:
                return extract_text_from_image_bytes_azure(f.read())
        else:
            raise ValueError(f"Format de fichier non supporté : {ext}")


# ----------- Nettoyage texte -----------
def clean_text(text: str) -> str:
    """
    Nettoie le texte en supprimant les espaces multiples et en normalisant.
    """
    if not text:
        return ""
    
    # Remplacer les sauts de ligne multiples par un espace
    text = re.sub(r'\n+', ' ', text)
    # Remplacer les espaces multiples par un seul espace
    text = re.sub(r'\s+', ' ', text)
    # Supprimer les espaces en début et fin
    text = text.strip()
    
    return text


# ----------- Découpage en chunks -----------
def split_into_chunks(text: str, max_tokens: int = 500, overlap_tokens: int = 50) -> List[str]:
    """
    Découpe un texte en chunks avec un chevauchement.
    
    Args:
        text: Texte à découper
        max_tokens: Nombre maximum de tokens par chunk
        overlap_tokens: Nombre de tokens de chevauchement entre les chunks
    
    Returns:
        Liste des chunks
    """
    # Nettoyer le texte d'abord
    text = clean_text(text)
    
    # Si le texte est vide
    if not text:
        return []
    
    # Diviser en phrases (ponctuation + espaces)
    sentences = re.split(r'(?<=[.!?])\s+', text)
    
    # Filtrer les phrases vides
    sentences = [s.strip() for s in sentences if s.strip()]
    
    # Si aucune phrase détectée ou texte très court, utiliser la méthode simple
    if len(sentences) <= 1 or len(text.split()) <= max_tokens:
        return split_into_chunks_simple(text, max_tokens, overlap_tokens)
    
    # Méthode par phrases (plus naturelle)
    chunks = []
    current_chunk = []
    current_length = 0
    
    for sentence in sentences:
        # Compter les mots dans la phrase (approximation des tokens)
        sentence_words = sentence.split()
        sentence_length = len(sentence_words)
        
        # Si la phrase seule dépasse max_tokens, la diviser
        if sentence_length > max_tokens:
            if current_chunk:  # Sauvegarder le chunk courant
                chunks.append(" ".join(current_chunk))
                current_chunk = []
                current_length = 0
            
            # Diviser la longue phrase
            sub_chunks = split_into_chunks_simple(sentence, max_tokens, overlap_tokens)
            chunks.extend(sub_chunks[:-1])  # Ajouter tous sauf le dernier
            
            # Le dernier sous-chunk devient le chunk courant
            if sub_chunks:
                current_chunk = [sub_chunks[-1]] if sub_chunks else []
                current_length = len(sub_chunks[-1].split()) if sub_chunks else 0
        
        # Si l'ajout de la phrase dépasse la limite
        elif current_length + sentence_length > max_tokens:
            if current_chunk:  # Éviter les chunks vides
                chunks.append(" ".join(current_chunk))
            
            # Commencer nouveau chunk avec chevauchement
            if overlap_tokens > 0 and current_chunk:
                # Prendre les dernières phrases pour le chevauchement
                overlap_words = []
                overlap_count = 0
                
                # Reconstruire le chevauchement depuis la fin du chunk précédent
                for prev_sentence in reversed(current_chunk):
                    prev_words = prev_sentence.split()
                    if overlap_count + len(prev_words) <= overlap_tokens:
                        overlap_words.insert(0, prev_sentence)
                        overlap_count += len(prev_words)
                    else:
                        # Prendre seulement une partie de la phrase si nécessaire
                        needed = overlap_tokens - overlap_count
                        if needed > 0:
                            partial_sentence = " ".join(prev_words[-needed:])
                            overlap_words.insert(0, partial_sentence)
                        break
                
                current_chunk = overlap_words + [sentence]
                current_length = overlap_count + sentence_length
            else:
                current_chunk = [sentence]
                current_length = sentence_length
        
        else:
            current_chunk.append(sentence)
            current_length += sentence_length
    
    # Ajouter le dernier chunk
    if current_chunk:
        chunks.append(" ".join(current_chunk))
    
    # Filtrer les chunks vides et s'assurer qu'ils ne sont pas trop longs
    final_chunks = []
    for chunk in chunks:
        chunk = chunk.strip()
        if chunk:
            words = chunk.split()
            if len(words) > max_tokens:
                # Redécouper les chunks trop longs
                sub_chunks = split_into_chunks_simple(chunk, max_tokens, overlap_tokens)
                final_chunks.extend(sub_chunks)
            else:
                final_chunks.append(chunk)
    
    return final_chunks


def split_into_chunks_simple(text: str, max_tokens: int = 500, overlap_tokens: int = 50) -> List[str]:
    """
    Version simplifiée pour un découpage plus agressif par mots.
    """
    text = clean_text(text)
    if not text:
        return []
    
    words = text.split()
    
    # Si le texte est plus court que max_tokens, retourner le texte entier
    if len(words) <= max_tokens:
        return [" ".join(words)] if words else []
    
    chunks = []
    
    i = 0
    while i < len(words):
        # Prendre un chunk de max_tokens
        chunk_end = min(i + max_tokens, len(words))
        chunk = words[i:chunk_end]
        
        if chunk:
            chunks.append(" ".join(chunk))
        
        # Avancer avec chevauchement
        i += max_tokens - overlap_tokens
        
        # Éviter la boucle infinie
        if i >= len(words):
            break
    
    return chunks


# ----------- Extraction de réponse, sources, résumé -----------
def extract_answer(response_text: str, summarize: bool = False) -> Tuple[str, list, str]:
    """
    Extrait la réponse, les sources et éventuellement un résumé du texte de réponse.
    """
    cleaned = response_text.strip()

    # Supprimer les préfixes de réponse courants
    prefixes = [
        r"^(Réponse\s*:|Réponse\s*-\s*|Answer\s*:|Answer\s*-\s*)",
        r"^La réponse est\s*[:\-]?",
        r"^Voici\s+la\s+r[eé]ponse\s*[:\-]?"
    ]
    for prefix in prefixes:
        cleaned = re.sub(prefix, "", cleaned, flags=re.IGNORECASE).strip()

    # Extraire les sources
    sources = re.findall(r"\[source:\s*(.*?)\]", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\[source:\s*.*?\]", "", cleaned, flags=re.IGNORECASE).strip()

    # Générer un résumé si demandé
    summary = ""
    if summarize:
        sentences = re.split(r'(?<=[.!?])\s+', cleaned)
        summary = " ".join(sentences[:2]) if len(sentences) > 2 else cleaned

    return cleaned, sources, summary


# ----------- Fonctions utilitaires supplémentaires -----------
def estimate_tokens(text: str) -> int:
    """
    Estime le nombre de tokens dans un texte (approximation: 1 token ≈ 4 caractères ou 0.75 mots)
    """
    if not text:
        return 0
    # Approximation simple: nombre de mots * 1.3
    return int(len(text.split()) * 1.3)


def get_file_info(file_path: str) -> dict:
    """
    Retourne des informations sur le fichier
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Fichier non trouvé: {file_path}")
    
    stat = os.stat(file_path)
    return {
        "filename": os.path.basename(file_path),
        "size": stat.st_size,
        "extension": os.path.splitext(file_path)[1].lower(),
        "modified": stat.st_mtime
    }


# Test de la fonction split_into_chunks
if __name__ == "__main__":
    # Test avec un texte long
    test_text = """
    Ceci est un texte de test. Il contient plusieurs phrases. 
    Chaque phrase devrait être séparée. Le but est de voir si le découpage fonctionne correctement.
    Nous allons répéter ce texte plusieurs fois pour avoir un contenu suffisamment long.
    Ceci est un texte de test. Il contient plusieurs phrases. 
    Chaque phrase devrait être séparée. Le but est de voir si le découpage fonctionne correctement.
    Nous allons répéter ce texte plusieurs fois pour avoir un contenu suffisamment long.
    Ceci est un texte de test. Il contient plusieurs phrases. 
    Chaque phrase devrait être séparée. Le but est de voir si le découpage fonctionne correctement.
    Nous allons répéter ce texte plusieurs fois pour avoir un contenu suffisamment long.
    """
    
    chunks = split_into_chunks(test_text, max_tokens=100, overlap_tokens=20)
    print(f"Nombre de chunks générés: {len(chunks)}")
    for i, chunk in enumerate(chunks):
        words = chunk.split()
        print(f"Chunk {i+1}: {len(words)} mots, {estimate_tokens(chunk)} tokens estimés")
        print(f"Preview: {chunk[:80]}...")
        print("---")